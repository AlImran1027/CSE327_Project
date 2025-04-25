import os
import requests
from pptx import Presentation
from pptx.util import Inches, Pt

def download_image(url, filename):
    response = requests.get(url)
    if response.status_code == 200:
        with open(filename, 'wb') as f:
            f.write(response.content)
        return filename
    else:
        print("Failed to download", url)
        return None

# Define image URLs (using Unsplash's source service for direct access)
url_slide4 = "https://source.unsplash.com/dE6TqH-i94Q/800x600"  # Worksheet image
url_slide5 = "https://source.unsplash.com/vZlTg_McCDo/800x600"   # Scaffolding image
url_slide7 = "https://source.unsplash.com/3Kv48NS4WUU/800x600"    # Final work image

# Download images
img_slide4 = download_image(url_slide4, "slide4.jpg")
img_slide5 = download_image(url_slide5, "slide5.jpg")
img_slide7 = download_image(url_slide7, "slide7.jpg")

# Create a new presentation
prs = Presentation()

# Slide 1: Cover
slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide layout
title_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
tf = title_box.text_frame
tf.text = "Analyzing a 5th Grade ELL Student’s Writing Development"
p = tf.add_paragraph()
p.text = "A Study on Discipline-Specific Writing in Literature"
p.font.size = Pt(24)
p = tf.add_paragraph()
p.text = "Your Name"
p = tf.add_paragraph()
p.text = "Date"

# Slide 2: Class Info
slide = prs.slides.add_slide(prs.slide_layouts[5])
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
title_box.text = "Class Info"
content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(4))
content = content_box.text_frame
content.text = (
    "Grade Level: 5th Grade\n"
    "School Location: [Insert Location]\n"
    "Type of School: [Public/Private/Charter]\n"
    "Type of Classroom: [Co-Taught or General Education/ESL Pull-Out]\n"
    "Content Area: English Language Arts (Literature)"
)

# Slide 3: Student Info
slide = prs.slides.add_slide(prs.slide_layouts[5])
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
title_box.text = "Student Info"
content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(4))
content = content_box.text_frame
content.text = (
    "Language Proficiency: Emerging Bilingual (Intermediate Level)\n"
    "Background: [General details about the student’s language background]\n"
    "Personality: Engaged, eager to learn, participates in discussions"
)

# Slide 4: Part 1 Observation Overview
slide = prs.slides.add_slide(prs.slide_layouts[5])
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
title_box.text = "Part 1 Observation Overview"
content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(5), Inches(4))
content = content_box.text_frame
content.text = (
    "Lesson Objective: Identify character traits in a novel using textual evidence\n\n"
    "Lesson Materials:\n"
    " - Character Trait Worksheet\n"
    " - Graphic Organizer\n"
    " - Selected excerpts from an age-appropriate novel (e.g., Charlotte’s Web)\n\n"
    "Student Task: Choose a character, assign a trait, and provide evidence from the text\n\n"
    "Scaffolding Strategies:\n"
    " - Sentence Starters\n"
    " - Word Bank with Character Traits\n"
    " - Teacher Modeled Example\n\n"
    "Picture: Example of a Character Trait Worksheet"
)
# Add image to slide 4
if img_slide4:
    left = Inches(6)
    top = Inches(1.2)
    slide.shapes.add_picture(img_slide4, left, top, width=Inches(3), height=Inches(2.25))

# Slide 5: Part 2 Observation Overview
slide = prs.slides.add_slide(prs.slide_layouts[5])
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
title_box.text = "Part 2 Observation Overview"
content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(5), Inches(4))
content = content_box.text_frame
content.text = (
    "Lesson Objective: Develop written responses explaining character traits\n\n"
    "Lesson Materials:\n"
    " - Response Sheet with Sentence Frames\n"
    " - Peer Discussion Prompts\n\n"
    "Student Task: Write a short paragraph explaining a character’s trait using evidence\n\n"
    "Scaffolding Strategies:\n"
    " - Partner Brainstorming\n"
    " - Guided Questions from Teacher\n"
    " - Example Sentences on the Board\n\n"
    "Picture: Example of Scaffolding Materials (Sentence Frames & Organizers)"
)
# Add image to slide 5
if img_slide5:
    left = Inches(6)
    top = Inches(1.2)
    slide.shapes.add_picture(img_slide5, left, top, width=Inches(3), height=Inches(2.25))

# Slide 6: Part 3 Observation Overview
slide = prs.slides.add_slide(prs.slide_layouts[5])
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
title_box.text = "Part 3 Observation Overview"
content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(4))
content = content_box.text_frame
content.text = (
    "Lesson Objective: Participate in a discussion on character traits\n\n"
    "Lesson Materials:\n"
    " - Discussion Questions\n"
    " - Student’s Written Response from Previous Lesson\n\n"
    "Student Task: Verbally share their analysis with peers\n\n"
    "Scaffolding Strategies:\n"
    " - Small Group Discussion First\n"
    " - Sentence Starters for Speaking\n"
    " - Teacher Feedback and Rephrasing"
)

# Slide 7: Conclusion
slide = prs.slides.add_slide(prs.slide_layouts[5])
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
title_box.text = "Conclusion"
content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(5), Inches(4))
content = content_box.text_frame
content.text = (
    "Student’s Final Work Product:\n"
    " - Picture of the student’s final written response\n\n"
    "Growth in Language Acquisition:\n"
    " - Began with single-word answers, developed into full sentences\n"
    " - Improved ability to cite evidence from the text\n"
    " - Used more complex sentence structures in discussion\n\n"
    "Final Thought: The student demonstrated significant growth in both written and verbal expression with the support of scaffolding techniques."
)
# Add image to slide 7
if img_slide7:
    left = Inches(6)
    top = Inches(1.2)
    slide.shapes.add_picture(img_slide7, left, top, width=Inches(3), height=Inches(2.25))

# Save the presentation to a file
output_filename = "ELL_Student_Presentation_Proper.pptx"
prs.save(output_filename)
print("Presentation created successfully and saved as", output_filename)

